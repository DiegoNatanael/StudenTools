# ppt_generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

def create_presentation_pptx(title_text: str, subtitle_text: str, agenda_items: list, features_items: list) -> bytes:
    """
    Generates a .pptx presentation based on provided data and returns its bytes.
    """
    prs = Presentation()

    # --- Slide 1: Title slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text

    title_tf = title.text_frame
    for paragraph in title_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(40)

    # --- Slide 2: Agenda (bullet points) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Agenda"

    body_tf = slide.shapes.placeholders[1].text_frame
    body_tf.clear()
    p = body_tf.paragraphs[0]
    p.text = "Introducción"
    p.level = 0

    for item in agenda_items:
        p = body_tf.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Two columns (text + visual placeholder) ---
    slide_layout = prs.slide_layouts[5] # A blank-like layout
    slide = prs.slides.add_slide(slide_layout)

    # Left column: Title + bullets
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.2), Inches(4.0))
    left_tf = left.text_frame
    left_tf.text = "Características clave"
    left_tf.paragraphs[0].font.size = Pt(28)

    for item in features_items:
        p = left_tf.add_paragraph()
        p.text = "• " + item
        p.level = 1
        p.font.size = Pt(18)

    # Right column: "visual" placeholder
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.0), Inches(4.0), Inches(3.0))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(232, 242, 255)
    line = shape.line
    line.width = Pt(1)
    shape.text = "Espacio para imagen o gráfico\nPuedes reemplazarlo con tu imagen"

    shape_tf = shape.text_frame
    shape_tf.paragraphs[0].font.size = Pt(14)
    shape_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slide 4: Contact / Close ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Cierre y contacto"

    body_tf = slide.shapes.placeholders[1].text_frame
    body_tf.clear()
    p = body_tf.paragraphs[0]
    p.text = "Gracias por tu atención"
    p.font.size = Pt(20)
    p = body_tf.add_paragraph()
    p.text = "Email: tu.email@ejemplo.com"
    p.font.size = Pt(14)
    p = body_tf.add_paragraph()
    p.text = "¿Quieres que personalice esto con tus colores, logo o contenido?"
    p.font.size = Pt(14)

    # Save to a BytesIO object instead of a file
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
